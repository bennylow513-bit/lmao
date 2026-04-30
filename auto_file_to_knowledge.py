import csv
import json
import time
from datetime import datetime
from pathlib import Path
from typing import Optional


from watchdog.events import FileSystemEventHandler
from watchdog.observers import Observer


print("SCRIPT STARTED")

UPLOAD_FOLDER = Path("uploads")
KNOWLEDGE_FILE = Path("knowledge.txt")
PROCESSED_FILE = Path("processed_files.json")

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".md",
    ".docx",
    ".csv",
    ".xlsx",
    ".pptx",
    ".json",
}


# =========================
# PROCESSED FILE TRACKING
# =========================

def load_processed_files() -> dict:
    if not PROCESSED_FILE.exists():
        return {}

    try:
        with open(PROCESSED_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def save_processed_files(data: dict) -> None:
    with open(PROCESSED_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def get_file_signature(file_path: Path) -> str:
    stat = file_path.stat()
    return f"{file_path.name}|{stat.st_size}|{int(stat.st_mtime)}"


def is_already_processed(file_path: Path) -> bool:
    processed = load_processed_files()
    signature = get_file_signature(file_path)
    return processed.get(file_path.name) == signature


def mark_as_processed(file_path: Path) -> None:
    processed = load_processed_files()
    processed[file_path.name] = get_file_signature(file_path)
    save_processed_files(processed)


# =========================
# WAIT UNTIL FILE IS READY
# =========================

def wait_until_file_is_ready(file_path: Path, checks: int = 5, delay: float = 1.0) -> bool:
    """
    This prevents reading the file while it is still being copied into the folder.
    """
    last_size = -1

    for _ in range(checks):
        if not file_path.exists():
            return False

        current_size = file_path.stat().st_size

        if current_size == last_size and current_size > 0:
            return True

        last_size = current_size
        time.sleep(delay)

    return file_path.exists() and file_path.stat().st_size > 0


# =========================
# TEXT CLEANING
# =========================

def clean_text(text: str) -> str:
    lines = []

    for line in text.splitlines():
        line = line.strip()

        if line:
            lines.append(line)

    return "\n".join(lines)


# =========================
# FILE EXTRACTORS
# =========================

def extract_pdf(file_path: Path) -> str:
    try:
        import fitz
    except ImportError:
        return "ERROR: pymupdf is not installed. Run: pip install pymupdf"

    print(f"Reading PDF: {file_path}")

    text_parts = []

    try:
        doc = fitz.open(file_path)

        for page_number, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()

            if text:
                text_parts.append(f"\n--- Page {page_number} ---\n{text}")

        doc.close()

    except Exception as e:
        return f"ERROR reading PDF: {e}"

    return "\n".join(text_parts).strip()


def extract_txt_or_md(file_path: Path) -> str:
    print(f"Reading text file: {file_path}")

    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="latin-1")
    except Exception as e:
        return f"ERROR reading text file: {e}"


def extract_docx(file_path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return "ERROR: python-docx is not installed. Run: pip install python-docx"

    print(f"Reading DOCX: {file_path}")

    try:
        document = Document(file_path)
        parts = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table_number, table in enumerate(document.tables, start=1):
            parts.append(f"\n--- Table {table_number} ---")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)

        return "\n".join(parts)

    except Exception as e:
        return f"ERROR reading DOCX: {e}"


def extract_csv(file_path: Path) -> str:
    print(f"Reading CSV: {file_path}")

    try:
        parts = []

        with open(file_path, "r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f)

            for row_number, row in enumerate(reader, start=1):
                row_text = " | ".join(str(cell).strip() for cell in row)
                parts.append(f"Row {row_number}: {row_text}")

        return "\n".join(parts)

    except Exception as e:
        return f"ERROR reading CSV: {e}"


def extract_xlsx(file_path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return "ERROR: openpyxl is not installed. Run: pip install openpyxl"

    print(f"Reading XLSX: {file_path}")

    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        parts = []

        for sheet in workbook.worksheets:
            parts.append(f"\n--- Sheet: {sheet.title} ---")

            for row in sheet.iter_rows(values_only=True):
                values = []

                for cell in row:
                    if cell is None:
                        values.append("")
                    else:
                        values.append(str(cell).strip())

                if any(values):
                    parts.append(" | ".join(values))

        return "\n".join(parts)

    except Exception as e:
        return f"ERROR reading XLSX: {e}"


def extract_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

    print(f"Reading PPTX: {file_path}")

    try:
        presentation = Presentation(file_path)
        parts = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts.append(f"\n--- Slide {slide_number} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        parts.append(text)

        return "\n".join(parts)

    except Exception as e:
        return f"ERROR reading PPTX: {e}"


def extract_json(file_path: Path) -> str:
    print(f"Reading JSON: {file_path}")

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        return json.dumps(data, ensure_ascii=False, indent=2)

    except Exception as e:
        return f"ERROR reading JSON: {e}"


def extract_text_from_file(file_path: Path) -> Optional[str]:
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_pdf(file_path)

    if suffix in {".txt", ".md"}:
        return extract_txt_or_md(file_path)

    if suffix == ".docx":
        return extract_docx(file_path)

    if suffix == ".csv":
        return extract_csv(file_path)

    if suffix == ".xlsx":
        return extract_xlsx(file_path)

    if suffix == ".pptx":
        return extract_pptx(file_path)

    if suffix == ".json":
        return extract_json(file_path)

    return None


# =========================
# ADD TO KNOWLEDGE
# =========================

def append_file_to_knowledge(file_path: Path) -> None:
    suffix = file_path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        print(f"Skipped unsupported file type: {file_path.name}")
        return

    if is_already_processed(file_path):
        print(f"Already processed and unchanged: {file_path.name}")
        return

    if not wait_until_file_is_ready(file_path):
        print(f"File is not ready: {file_path.name}")
        return

    print(f"Processing file: {file_path.name}")

    raw_text = extract_text_from_file(file_path)

    if not raw_text:
        print(f"No text extracted from {file_path.name}")
        return

    cleaned_text = clean_text(raw_text)

    if not cleaned_text:
        print(f"No usable text found in {file_path.name}")
        return

    today = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    section = f"""

==================================================
AUTO-IMPORTED FILE: {file_path.name}
File type: {suffix}
Imported on: {today}
==================================================

Important:
- This section was auto-extracted from an uploaded file.
- Check and clean this section before final submission.
- If any information conflicts with earlier confirmed knowledge, update the correct section manually.
- Do not treat unclear or incomplete information as confirmed.

{cleaned_text}

==================================================
END OF AUTO-IMPORTED FILE: {file_path.name}
==================================================

"""

    with open(KNOWLEDGE_FILE, "a", encoding="utf-8") as f:
        f.write(section)

    mark_as_processed(file_path)

    print(f"Done. Added {file_path.name} into knowledge.txt")


# =========================
# WATCHER
# =========================

class UploadHandler(FileSystemEventHandler):
    def on_created(self, event):
        if event.is_directory:
            return

        file_path = Path(event.src_path)

        if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            print(f"New file detected: {file_path.name}")
            append_file_to_knowledge(file_path)

    def on_modified(self, event):
        if event.is_directory:
            return

        file_path = Path(event.src_path)

        if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            print(f"File modified: {file_path.name}")
            append_file_to_knowledge(file_path)


def process_existing_files():
    UPLOAD_FOLDER.mkdir(exist_ok=True)

    files = [
        file_path
        for file_path in UPLOAD_FOLDER.iterdir()
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS
    ]

    if not files:
        print("No supported files found yet in uploads/ folder.")
        return

    print(f"Found {len(files)} supported file(s).")

    for file_path in files:
        append_file_to_knowledge(file_path)


def watch_upload_folder():
    UPLOAD_FOLDER.mkdir(exist_ok=True)

    print("Watching uploads/ folder...")
    print("Supported files:")
    print(", ".join(sorted(SUPPORTED_EXTENSIONS)))
    print("Drop a supported file into uploads/ and it will be added to knowledge.txt.")
    print("Press CTRL + C to stop.")

    process_existing_files()

    event_handler = UploadHandler()
    observer = Observer()
    observer.schedule(event_handler, str(UPLOAD_FOLDER), recursive=False)
    observer.start()

    try:
        while True:
            time.sleep(1)

    except KeyboardInterrupt:
        print("Stopping watcher...")
        observer.stop()

    observer.join()


if __name__ == "__main__":
    watch_upload_folder()